"""Build an editable PPTX and a matching offline preview from project assets."""
from pathlib import Path
import sys, json, math, html, shutil, zipfile

HERE = Path(__file__).resolve().parent
REPO = HERE.parent
TOOLS = REPO.parent / '.tools'
sys.path[:0] = [str(TOOLS / 'presentation'), str(TOOLS / 'title-preview')]
from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from playwright.sync_api import sync_playwright

ASSETS = HERE / 'media'
ASSETS.mkdir(exist_ok=True)
W, H = 1600, 900
INK, MUTED, GREEN, BLUE = '#203a36', '#627571', '#36795a', '#427f9c'
FONT = 'Microsoft YaHei'
slides = []

def rgb(c):
    return tuple(int(c[i:i+2], 16) for i in (1, 3, 5))

def raster_assets():
    with sync_playwright() as p:
        browser = p.chromium.launch(executable_path=r'C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe', headless=True, args=['--allow-file-access-from-files'])
        page = browser.new_page(viewport={'width': 1520, 'height': 1280}, device_scale_factor=1)
        for name in ['library', 'hongshan', 'gate', 'redhall']:
            src = REPO / 'assets/buildings/isometric' / f'{name}-0.svg'
            svg = src.read_text(encoding='utf-8')
            page.set_content(f'<html><style>body{{margin:0;background:transparent}}svg{{width:1520px;height:1280px}}</style><body>{svg}</body></html>')
            page.evaluate('document.fonts.ready')
            dest = ASSETS / f'{name}.png'
            page.screenshot(path=str(dest), omit_background=True)
            im = Image.open(dest)
            im.crop(im.getbbox()).save(dest)
        browser.close()
    for name in ['library', 'hongshan', 'gate', 'redhall']:
        shutil.copy2(REPO / f'assets/buildings/references/{name}.jpg', ASSETS / f'photo-{name}.jpg')

def render_terrain():
    models = json.loads((HERE / 'terrain_meshes.json').read_text(encoding='utf-8'))
    az = el = math.radians(42)
    sa, ca, se, ce = math.sin(az), math.cos(az), math.sin(el), math.cos(el)
    def project(v):
        x, y, z = v
        return ca*x-sa*z, se*(sa*x+ca*z)-ce*y, ce*(sa*x+ca*z)+se*y
    for name, faces in models.items():
        polygons = []
        for tri, col in faces:
            if col[3] < .26:
                continue
            pts = [project(v) for v in tri]
            a, b, c = tri
            u, v = [b[i]-a[i] for i in range(3)], [c[i]-a[i] for i in range(3)]
            n = [u[1]*v[2]-u[2]*v[1], u[2]*v[0]-u[0]*v[2], u[0]*v[1]-u[1]*v[0]]
            length = math.sqrt(sum(k*k for k in n)) or 1
            light = .83 + .17*abs((n[0]*.3+n[1]*.9+n[2]*.15)/length)
            color = tuple(min(255, round(k*255*light)) for k in col[:3]) + (255,)
            scale, cy = (5000,780) if name.startswith('flower') else (540,530)
            polygons.append((sum(p[2] for p in pts)/3, [(480+p[0]*scale, cy+p[1]*scale) for p in pts], color))
        canvas = Image.new('RGBA', (960, 1000))
        draw = ImageDraw.Draw(canvas)
        for _, pts, color in sorted(polygons, key=lambda x:x[0]):
            draw.polygon(pts, fill=color)
        canvas.crop(canvas.getbbox()).save(ASSETS / f'terrain-{name}.png')

def backdrop():
    image = Image.new('RGB', (W, H))
    d = ImageDraw.Draw(image)
    stops = [(0, '#c9e6f1'), (.64, '#edf4f5'), (1, '#fcfdfd')]
    for y in range(H):
        t = y/(H-1)
        a,b = (stops[0], stops[1]) if t < .64 else (stops[1], stops[2])
        v=(t-a[0])/(b[0]-a[0])
        color=tuple(round(x*(1-v)+z*v) for x,z in zip(rgb(a[1]),rgb(b[1])))
        d.line((0,y,W,y),fill=color)
    image.save(ASSETS/'background.png')

def image_on(canvas, path, box):
    im = Image.open(path).convert('RGBA')
    x,y,w,h=box
    im.thumbnail((int(w),int(h)), Image.Resampling.LANCZOS)
    canvas.alpha_composite(im, (int(x+(w-im.width)/2),int(y+h-im.height)))

def city_scene():
    canvas=Image.new('RGBA',(1600,750))
    cells=[]
    for x in range(7):
        for z in range(5):
            typ = 1 if x == 3 else (5 if x in (0,6) and z in (0,4) else [0,2,0,3][(x*3+z)%4])
            if (x,z) in [(1,2),(4,1),(5,3),(2,4)]: typ=4
            cells.append((x+z,x,z,typ))
    for _,x,z,typ in sorted(cells):
        sx=790+(x-z-1)*114
        sy=240+(x+z)*49
        image_on(canvas,ASSETS/f'terrain-{typ}.png',(sx-125,sy-150,250,245))
        if typ in (0,2,3):
            for index in range(3):
                image_on(canvas,ASSETS/f'terrain-flower{(x+z)%4}.png',(sx-45+index*30,sy+8+(index%2)*8,23,42))
        name={(1,2):'library',(4,1):'hongshan',(5,3):'redhall',(2,4):'gate'}.get((x,z))
        if name:
            image_on(canvas,ASSETS/f'{name}.png',(sx-180,sy-230,360,320))
    canvas.save(ASSETS/'city.png')

def new(title, kicker, notes, seconds=25):
    s={'title':title,'notes':notes,'seconds':seconds,'items':[]}
    slides.append(s)
    img(s,'background.png',0,0,W,H,fit='stretch')
    rect(s,48,34,1504,48,'#e7eeee',stroke='#ffffff')
    text(s,'花满洪山  /  HONGSHAN IN BLOOM',68,44,900,26,17,MUTED)
    text(s,kicker,1130,44,396,26,17,MUTED,align='right')
    if title: text(s,title,70,119,1460,92,54,INK,bold=True)
    text(s,'洪山区多样性主题游戏',70,855,700,24,16,MUTED)
    text(s,f'{len(slides):02d}',1430,850,100,30,20,MUTED,align='right')
    return s

def text(s, value,x,y,w,h,size=26,color=INK,bold=False,align='left',font=FONT):
    s['items'].append(dict(kind='text',text=value,x=x,y=y,w=w,h=h,size=size,color=color,bold=bold,align=align,font=font))

def rect(s,x,y,w,h,color='#eef2f1',stroke=None,radius=0):
    s['items'].append(dict(kind='rect',x=x,y=y,w=w,h=h,color=color,stroke=stroke,radius=radius))

def img(s,name,x,y,w,h,fit='contain'):
    s['items'].append(dict(kind='image',src=name,x=x,y=y,w=w,h=h,fit=fit))

def rule(s, x, y, number, title, body, w=440):
    text(s,number,x,y,100,52,42,GREEN,bold=True)
    text(s,title,x,y+69,w,46,31,INK,bold=True)
    text(s,body,x,y+124,w,115,24,MUTED)

def card(s,x,y,kind,title,subtitle,model, accent,level='I'):
    width,height=275,353
    # Match _draw_card_component: offset shadow, accent rim, light-backed gradient.
    rect(s,x+8,y+11,width,height,'#b4bfbe')
    rect(s,x,y,width,height,accent)
    start=rgb(accent); target=rgb('#f0d0b0' if kind=='develop' else '#d0e4f4' if kind=='weather' else '#fffff7')
    for row in range(64):
        t=row/63
        a,b,u=(start,target,t/.45) if t<.45 else (target,(255,255,247),(t-.45)/.55)
        color='#'+''.join(f'{round(a[i]*(1-u)+b[i]*u):02x}' for i in range(3))
        rect(s,x+5,y+5+row*(height-10)/64,width-10,(height-10)/64+1,color)
    text(s,title,x+20,y+19,width-40,45,30,INK,bold=True)
    text(s,level,x+210,y+22,40,32,19,INK,align='right')
    if model: img(s,model,x+25,y+85,width-50,174)
    if kind=='weather':
        text(s,'☁',x+35,y+73,205,134,98,'#619ab5',align='center',font='Segoe UI Symbol')
        text(s,'雨  /  风  /  旱  /  虹',x+21,y+225,width-42,35,20,BLUE,align='center')
    text(s,subtitle,x+19,y+286,width-38,45,21,INK,align='center')

def make_slides():
    s=new('', '作品介绍', '大家好，这是花满洪山，一款把洪山的山水、高校与科创建筑带上棋盘的轻策略游戏。玩家不只是参观地标，而是通过播种、开发、连路和天气，让自己的花色在城市里生长。',15)
    text(s,'花满洪山',70,119,970,154,112,INK,True,font='KaiTi')
    text(s,'让一座大学之城，在棋盘上开花。',77,279,1250,59,35,INK)
    text(s,'2.5D 地块  ·  卡牌策略  ·  城市文化',80,351,1140,38,23,MUTED)
    img(s,'city.png',125,323,1390,505)
    rect(s,1165,174,365,154,'#edf2f2',stroke='#ffffff')
    text(s,'山水 / 人文 / 创新',1188,194,325,35,26,GREEN,True)
    text(s,'洪山区多样性主题\nHongshan in Bloom',1188,244,325,66,21,MUTED)

    s=new('把城市的多样性，变成可玩的关系','01 / 设计出发点','我们用三个维度理解洪山的多样性。山水生态被转化为不同的地块与生长条件；高校文化由图书馆和校园建筑承载；科创空间则成为能影响周边的建筑节点。它们不是互不相关的装饰，而是在同一张棋盘上彼此影响。这里的多样性是城市文化与景观的设计表达，不是现实生态普查。',23)
    for x,model,no,heading,body in [(76,'terrain-1.png','01','山水生态','森林、草地与水域\n共同塑造生长环境'),(589,'library.png','02','大学文化','图书馆与校园地标\n承载学习与城市记忆'),(1102,'hongshan.png','03','科创空间','科创楼宇进入棋盘\n转化为周边增益节点')]:
        img(s,model,x+15,222,410,275)
        rule(s,x+10,537,no,heading,body,420)

    s=new('一回合：抽 3 张，出牌，再生长','02 / 玩法总览','游戏支持一到四人，竞争玩法建议两到四人。每位玩家选择一种花色，开局有五张播种卡。回合开始，从开发、道路、天气三个公共卡堆中任意组合抽三张，然后自由出牌。结束回合后，依次结算地块变化、花朵增值和扩散，再轮到下一位。每人行动十次后，比花朵总数；平局再比占据的植物地块数。',28)
    for x,n,title,body in [(78,'01','抽牌','三个公共卡堆\n任意组合抽取 3 张'),(592,'02','出牌','可使用任意数量手牌\n也可以直接结束回合'),(1106,'03','结算','地块变化 → 花朵增值\n→ 花朵扩散 → 下一位')]:
        rule(s,x,245,n,title,body,430)
    rect(s,72,552,1456,238,'#edf2f1',stroke='#ffffff')
    for i,(color,label) in enumerate([('#d83232','虞美人'),('#8c4edb','风铃草'),('#e0802f','万寿菊'),('#1c497d','鸢尾花')]):
        rect(s,105+i*225,584,12,43,color)
        text(s,label,133+i*225,591,180,40,28,color,True)
    text(s,'每人 5 张起始播种卡',104,666,600,45,31,INK,True)
    text(s,'每人行动 10 次后，花朵最多者获胜。',104,721,1310,39,27,MUTED)

    s=new('四类卡牌，四种改变棋盘的方式','03 / 卡牌操作','播种卡把十到五十朵自己的花放到未满的植物地块。开发卡把山体变成一组随机地形，建筑开发卡则把缺口变成指定建筑。道路卡连续拖动连接两到四格。天气卡改变全局：极端天气持续三回合并降低生长率；彩虹清除极端天气并增强生长与扩散。先把卡拖出手牌，再预览或放置，不会因为点一下就误用。',33)
    data=[('seed','播种卡','增加 10–50 朵花','terrain-0.png','#4f9c62','I–V'),('develop','开发卡','开发山体 / 建造建筑','terrain-5.png','#c47a42','I–IV'),('road','道路卡','连续拖动连接 2–4 格','terrain-road.png','#ae8b55','I–III'),('weather','天气卡','改变本轮生长环境',None,'#568eb0','')]
    for i,values in enumerate(data):
        x=88+i*386
        card(s,x,230,*values)
        desc=['只能播种在未满的\n森林、草地、荒漠','先看随机地形\n再旋转、选择位置','上下左右连续连接\n可以转弯和接入旧路','极端天气持续 3 回合\n彩虹清除极端天气'][i]
        text(s,desc,x-7,630,310,93,25,MUTED,align='center')
    for i in range(3): img(s,'terrain-flower0.png',165+i*30,369+(i%2)*12,44,80)
    text(s,'拖出手牌 → 悬停预览 → 点击放置    /    Q、E 旋转    /    Esc 或右键取消',78,778,1460,41,24,INK)

    s=new('同一片土地，不同的生长选择','04 / 地块与增益','森林容量八十、生长率百分之四十，草地是五十和百分之三十，荒漠是三十和百分之二十。环境不是固定的：植物地块会受邻居和天气影响，逐级升级或退化。水域能使周围八格的植物升级概率翻倍，建筑则让同样范围内的花朵容量翻倍。多个同类增益不叠加。因此，抢好的位置和营造好的邻居同样重要。',27)
    for x,kind,title,cap,growth in [(85,2,'森林','80','40%'),(427,0,'草地','50','30%'),(769,3,'荒漠','30','20%')]:
        img(s,f'terrain-{kind}.png',x,222,290,262)
        text(s,title,x,508,290,48,34,INK,True,align='center')
        text(s,f'容积 {cap}  /  生长率 {growth}',x-4,571,300,43,24,MUTED,align='center')
    rect(s,1121,222,401,475,'#edf2f1',stroke='#ffffff')
    img(s,'terrain-1.png',1146,249,145,145)
    text(s,'水域',1300,266,185,42,28,BLUE,True)
    text(s,'升级概率 ×2',1298,321,200,40,24,INK)
    img(s,'library.png',1145,444,160,145)
    text(s,'建筑',1300,461,185,42,28,GREEN,True)
    text(s,'花朵容积 ×2',1298,516,200,40,24,INK)
    text(s,'均影响周围 8 格\n多个同类增益不叠加',1147,616,354,62,21,MUTED)
    text(s,'荒漠 ⇄ 草地 ⇄ 森林',88,685,935,66,37,GREEN,True,align='center')
    text(s,'受邻居与天气影响，每次只改变一级。',90,758,940,39,25,MUTED,align='center')

    s=new('先规划形状，再把道路连成环','05 / 放置规则','开发不是逐个点格子，而是放置一到四格的组合。整组必须覆盖山体，并至少有一格与已开发区正交相邻。被开发区完全包围的内部山体会成为缺口，之后才能用建筑卡建造。建成建筑会获得一级播种卡。道路的开放端点全部连接并形成闭合区域后，满足花朵条件的玩家获得播种卡，每条闭环只奖励一次。',28)
    text(s,'开发山体',80,232,610,54,35,INK,True)
    for x,y in [(145,365),(270,415),(145,465)]: img(s,'terrain-5.png',x,y,220,205)
    text(s,'1–4 格组合，可旋转',93,302,560,43,27,GREEN)
    text(s,'全部覆盖山体\n至少一格正交邻接已开发区',80,681,655,91,27,MUTED)
    text(s,'填补缺口 / 闭合道路',798,232,715,54,35,INK,True)
    img(s,'terrain-6.png',813,330,227,201)
    text(s,'→',1060,390,96,70,50,GREEN)
    img(s,'library.png',1168,310,315,250)
    text(s,'缺口上建建筑，获得 1 级播种卡',810,576,691,43,28,GREEN,True)
    text(s,'道路没有开放端点并形成闭合区域\n满足条件的玩家获卡，每条闭环只奖励一次',808,663,690,104,25,MUTED)
    text(s,'内部山体被非可开发地块完全包围后，自动识别为缺口。',80,802,1430,32,21,MUTED)

    s=new('从真实地标，到棋盘里的城市节点','06 / 洪山建筑','这里展示两座洪山建筑的实景与游戏模型。左边是武汉理工大学南湖图书馆，我们保留阶梯体量和横向层次，让知识空间成为一级建筑。右边是洪山科创大厦，保留塔楼、玻璃幕墙和裙房的辨识特征，作为两格的二级建筑。它们在游戏中提供共同的周边增益，是把大学与产业的城市联系转化为玩法的设计表达。',28)
    building_pair(s,78,'library','武汉理工大学南湖图书馆','I 级 · 1×1 缺口','阶梯体量 / 知识与公共空间','照片：gooood / 华南理工大学建筑设计研究院')
    building_pair(s,842,'hongshan','洪山科创大厦','II 级 · 1×2 缺口','塔楼与裙房 / 科技创新空间','照片：项目招商资料；区校合作背景见武汉市科技局')

    s=new('校园记忆，也可以成为一张建筑卡','07 / 高校文化','华中科技大学的建筑卡参考校园红色山墙建筑，照片和模型都能看到这一轮廓。武汉大学牌坊也进入了游戏，作为武汉周边高校文化的联动元素。需要明确，武汉大学本部位于武昌区，所以我们没有把它当作洪山区行政范围内的地标。模型采用简化色块与两档阴影，追求棋盘里的清晰识别，而不是测绘级复原。',26)
    building_pair(s,78,'redhall','华中科技大学','II 级 · 1×2 缺口','红色山墙 / 校园文化记忆','照片：SICAS 校园资料；模型为风格化校园意象')
    building_pair(s,842,'gate','武汉大学牌坊','II 级 · 1×2 缺口','牌坊轮廓 / 周边高校文化联动','照片：新浪校园图文；本部地址属武昌区，非洪山区')

    s=new('让玩家记住的，不止是建筑的名字','08 / 设计价值','花满洪山希望玩家通过一次次选择，认识城市中不同空间的关系。山水地块让人理解环境影响，建筑卡让校园和科创地标可识别，道路和天气则带来取舍。最终，玩家留下的不只是一张高分棋盘，也是对洪山的生态、人文与创新气质的一次主动探索。谢谢大家。',18)
    img(s,'city.png',94,236,980,520)
    rule(s,1110,250,'01','看见地方','以真实地标建立辨识',410)
    text(s,'02  理解关系',1110,541,410,48,31,GREEN,True)
    text(s,'把空间影响转化为策略',1110,603,415,43,25,MUTED)
    text(s,'一城山水，四季花开。',81,760,1380,67,44,INK,True)

    s=new('资料来源与版本说明','附录 / 不计入讲解','本页备查，不纳入三到四分钟讲解。照片仅用于建筑造型对照与项目介绍，权利归原作者，未确认公开再分发许可。公开展播或商业使用前应确认授权或替换为自摄照片。PPT采用当前拉取版本8e40b2d；README个别奖励和旧容量描述与代码不一致，数值及奖励说明以当前代码为准。',0)
    refs=[('01  游戏规则与模型','README.md、scripts/main3d.gd；版本 8e40b2d。地块为源码网格离线投影，建筑直接取现有 SVG。',''),('02  图书馆照片 / 项目信息','gooood：武汉理工大学南湖校区图书馆；设计方供稿，摄影：邵峰。','https://www.gooood.cn/library-of-south-lake-campus-of-wuhan-university-of-technology-china-by-architectural-design-and-research-institute-of-scut.htm'),('03  洪山科创大厦照片 / 地域依据','项目招商资料；武汉市科学技术局《工作动态》（2024-06-18）。','https://kjj.wuhan.gov.cn/xwzx_8/gzdt/202406/t20240618_2417645.html'),('04  高校照片 / 地域核实','华中科技大学：SICAS 校园资料；武汉大学牌坊：新浪校园图文。','https://www.sicas.cn/School/183/Contents/110815105210771.shtml'),('05  地理口径与图片权利','武大本部属武昌区；照片权利归原作者，公开展播前需确认授权。完整链接见配套来源文档。','https://service.whu.edu.cn/')]
    for i,(head,body,url) in enumerate(refs):
        y=232+i*111
        text(s,head,80,y,1400,36,27,GREEN,True)
        text(s,body,80,y+44,1400,40,21,MUTED)
        if url: s['items'][-1]['url']=url

def building_pair(s,x,name,title,level,detail,source):
    text(s,title,x,223,680,50,32,INK,True)
    text(s,level,x,282,680,35,23,GREEN)
    img(s,f'photo-{name}.jpg',x,342,396,277,fit='cover')
    img(s,f'{name}.png',x+405,329,298,303)
    text(s,'实地照片',x,641,370,34,20,MUTED)
    text(s,'游戏内建筑模型',x+410,641,280,34,20,MUTED)
    text(s,detail,x,708,705,45,26,INK)
    text(s,source,x,774,700,54,16,MUTED)

def actual_box(o):
    x,y,w,h=[o[k] for k in ('x','y','w','h')]
    if o['kind']=='image' and o['fit']=='contain':
        iw,ih=Image.open(ASSETS/o['src']).size
        f=min(w/iw,h/ih)
        nw,nh=iw*f,ih*f
        return x+(w-nw)/2,y+(h-nh)/2,nw,nh
    return x,y,w,h

def pptx_export():
    p=Presentation();p.slide_width=Inches(16);p.slide_height=Inches(9)
    p.core_properties.title='花满洪山 | 洪山区多样性主题游戏'
    p.core_properties.subject='玩法、规则与城市建筑；约3分46秒'
    p.core_properties.author='花满洪山项目组'
    for data in slides:
        slide=p.slides.add_slide(p.slide_layouts[6])
        for o in data['items']:
            x,y,w,h=actual_box(o)
            box=[Inches(v/100) for v in (x,y,w,h)]
            if o['kind']=='image':
                path=ASSETS/o['src']
                if o['fit']=='cover':
                    picture=slide.shapes.add_picture(str(path),*box)
                    iw,ih=Image.open(path).size
                    ratio=(iw/ih)/(w/h)
                    if ratio>1: picture.crop_left=picture.crop_right=(1-1/ratio)/2
                    else: picture.crop_top=picture.crop_bottom=(1-ratio)/2
                else: slide.shapes.add_picture(str(path),*box)
            elif o['kind']=='rect':
                sh=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,*box)
                sh.fill.solid();sh.fill.fore_color.rgb=RGBColor.from_string(o['color'][1:])
                if o['stroke']:
                    sh.line.color.rgb=RGBColor.from_string(o['stroke'][1:]);sh.line.width=Pt(.7)
                else:sh.line.fill.background()
            else:
                sh=slide.shapes.add_textbox(*box)
                tf=sh.text_frame;tf.clear();tf.word_wrap=True
                tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
                for i,line in enumerate(o['text'].split('\n')):
                    para=tf.paragraphs[0] if i==0 else tf.add_paragraph()
                    para.alignment={'left':1,'center':2,'right':3}[o['align']]
                    para.space_before=Pt(0);para.space_after=Pt(0);para.line_spacing=1.25
                    r=para.add_run();r.text=line
                    r.font.name=o['font'];r.font.size=Pt(o['size']*.72);r.font.bold=o['bold']
                    r.font.color.rgb=RGBColor.from_string(o['color'][1:])
                    ea=OxmlElement('a:ea');ea.set('typeface',o['font']);r._r.get_or_add_rPr().append(ea)
                    if o.get('url'):r.hyperlink.address=o['url']
        slide.notes_slide.notes_text_frame.text=f"建议讲解：{data['seconds']} 秒\n\n{data['notes']}"
        transition=OxmlElement('p:transition');transition.set('spd','med');transition.append(OxmlElement('p:fade'))
        slide._element.insert(2,transition)
    path=HERE/'花满洪山_洪山区多样性_展示版.pptx'
    p.save(path)
    check=Presentation(path)
    assert len(check.slides)==10
    for i,slide in enumerate(check.slides):
        for sh in slide.shapes:
            assert sh.left>=0 and sh.top>=0 and sh.left+sh.width<=p.slide_width+100 and sh.top+sh.height<=p.slide_height+100,(i,sh.name)
    print('PPTX:',path,'Slides:',len(check.slides),'Seconds:',sum(s['seconds'] for s in slides))

def html_export():
    result=['<!doctype html><meta charset="utf-8"><title>花满洪山 · 展示稿</title><style>*{box-sizing:border-box;letter-spacing:0}body{margin:0;background:#d7e2e6;font-family:"Microsoft YaHei",sans-serif}.slide{position:relative;width:1600px;height:900px;overflow:hidden;margin:24px auto;break-after:page}.el{position:absolute;margin:0}.txt{white-space:pre-wrap;line-height:1.25;overflow:visible}@media print{@page{size:1600px 900px;margin:0}body{background:white}.slide{margin:0;width:1600px;height:900px}}</style>']
    for i,s in enumerate(slides):
        result.append(f'<section class="slide" id="slide-{i+1}">')
        for o in s['items']:
            x,y,w,h=actual_box(o)
            style=f'left:{x}px;top:{y}px;width:{w}px;height:{h}px;'
            if o['kind']=='image':
                result.append(f'<img class="el" style="{style}object-fit:{"cover" if o["fit"]=="cover" else "fill"}" src="media/{o["src"]}">')
            elif o['kind']=='rect':
                result.append(f'<div class="el" style="{style}background:{o["color"]};border:{"1px solid "+o["stroke"] if o["stroke"] else "none"}"></div>')
            else:
                result.append(f'<div class="el txt" style="{style}font-family:&quot;{o["font"]}&quot;,sans-serif;font-size:{o["size"]}px;color:{o["color"]};font-weight:{700 if o["bold"] else 400};text-align:{o["align"]}">{html.escape(o["text"])}</div>')
        result.append('</section>')
    (HERE/'预览.html').write_text(''.join(result),encoding='utf-8')
    with sync_playwright() as p:
        browser=p.chromium.launch(executable_path=r'C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe',headless=True,args=['--allow-file-access-from-files'])
        page=browser.new_page(viewport={'width':1648,'height':948},device_scale_factor=1)
        page.goto((HERE/'预览.html').as_uri())
        page.wait_for_function('Array.from(document.images).every(i=>i.complete&&i.naturalWidth>0)')
        errors=page.evaluate('''() => [...document.querySelectorAll('.txt')].filter(e=>e.scrollHeight>e.clientHeight+3||e.scrollWidth>e.clientWidth+3).map(e=>e.textContent)''')
        if errors: print('OVERFLOW:',errors)
        assert not errors,errors
        for i in range(len(slides)):
            page.locator(f'#slide-{i+1}').screenshot(path=str(HERE/f'slide-{i+1:02}.png'))
        page.pdf(path=str(HERE/'花满洪山_预览.pdf'),print_background=True,prefer_css_page_size=True)
        browser.close()
    thumbs=Image.new('RGB',(1280,1800),'#d7e2e6')
    for i in range(10):
        im=Image.open(HERE/f'slide-{i+1:02}.png');im.thumbnail((624,351))
        thumbs.paste(im,(8+(i%2)*640,8+(i//2)*360))
    thumbs.save(HERE/'总览.jpg',quality=92)
    script=['# 花满洪山 · 讲解稿\n','正文9页，合计约226秒。第10页为备查，不计入讲解。\n']
    for i,s in enumerate(slides): script.extend([f'\n## {i+1:02} {s["title"] or "花满洪山"}（{s["seconds"]}秒）\n',s['notes']+'\n'])
    (HERE/'讲解稿.md').write_text('\n'.join(script),encoding='utf-8')

if __name__=='__main__':
    raster_assets()
    render_terrain()
    backdrop()
    city_scene()
    make_slides()
    pptx_export()
    html_export()
